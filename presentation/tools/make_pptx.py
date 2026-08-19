#!/usr/bin/env python3
"""Pixel-faithful PowerPoint export, derived from the master. Never hand-edited.

Each slide is the master's own 1920x1080 rendering (tools/render_check.py
screenshots) placed full-bleed, with the real speaker-notes text — parsed out
of the master's <aside class="notes"> — in the PowerPoint notes field.
Regenerate after any master change; an export edited by hand is a second deck.
"""
import html as htmllib
import io, os, re, sys

HERE = os.path.dirname(os.path.abspath(__file__))
PRES = os.path.dirname(HERE)
MASTER = os.path.join(PRES, "how-to-use-the-skills.html")
SLIDES = os.path.join(PRES, "exports", "slides")
OUT = os.path.join(PRES, "exports", "how-to-use-the-skills.pptx")

def notes_from_master():
    doc = io.open(MASTER, encoding="utf-8").read()
    # HTML comments first: the master's header comment names <aside class="notes">
    # in prose, and a raw scan would start slide 1's note there and swallow the
    # comment tail plus the whole embedded-font <style> block.
    doc = re.sub(r"<!--.*?-->", "", doc, flags=re.S)
    notes = re.findall(r'<aside class="notes">(.*?)</aside>', doc, re.S)
    return [re.sub(r"\s+", " ", htmllib.unescape(re.sub(r"<[^>]+>", "", n))).strip()
            for n in notes]

def main():
    import json
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE

    notes = notes_from_master()
    pngs = sorted(f for f in os.listdir(SLIDES) if f.endswith(".png"))
    if len(notes) != len(pngs):
        print("FAIL: %d notes but %d slide images" % (len(notes), len(pngs)))
        return 1
    links_path = os.path.join(PRES, "exports", "links.json")
    links = json.load(open(links_path)) if os.path.exists(links_path) else {}

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    px = prs.slide_width / 1920  # EMU per master pixel
    blank = prs.slide_layouts[6]
    n_hot = 0
    for i, (png, note) in enumerate(zip(pngs, notes), 1):
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(os.path.join(SLIDES, png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
        # invisible hotspots where the master renders its links (links.json
        # is measured off the built master by render_check.py)
        for a in links.get(str(i), []):
            shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Emu(int(a["x"] * px)), Emu(int(a["y"] * px)),
                                     Emu(int(a["w"] * px)), Emu(int(a["h"] * px)))
            shp.fill.background()
            shp.line.fill.background()
            shp.click_action.hyperlink.address = a["href"]
            n_hot += 1
        s.notes_slide.notes_text_frame.text = note
    prs.save(OUT)
    print("wrote %s (%d slides, notes embedded, %d link hotspots)" % (OUT, len(pngs), n_hot))
    return 0

if __name__ == "__main__":
    sys.exit(main())
