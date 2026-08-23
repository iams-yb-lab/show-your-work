#!/usr/bin/env python3
"""Mechanical layout check for a PowerPoint rebuild — the PPTX's own mechcheck.

The HTML master is checked in the browser, which measures real text. A PPTX
cannot be measured that way, so wrapped-text height is ESTIMATED, pessimistically,
from run sizes and box width. Three tests per slide:

  1. bounds     — no shape, and no shape's estimated text, leaves the canvas
  2. collisions — no picture overlaps another picture or a shape
  3. font floor — no run below the deck's font floor

Exits non-zero on any finding, so a bad rebuild cannot pass silently. Slide 1 is
exempt from the collision test: text over a full-bleed title image is intended.
"""
import math, os, sys
from pptx import Presentation
from deckcfg import load

PICTURE = 13


def main():
    cfg = load()
    path = cfg.export(cfg.pptx_name)
    if not os.path.exists(path):
        print("no such file: %s" % path, file=sys.stderr)
        return 2
    prs = Presentation(path)
    CW, CH = cfg.canvas
    PX = prs.slide_width / float(CW)
    FLOOR_PT = cfg.font_floor_px * 960.0 / CW      # the floor in points (960 pt = 13.333 in canvas)

    def px(v):
        return v / PX

    def est_text_h(sh):
        """pessimistic wrapped-text height, in canvas px"""
        if not sh.has_text_frame:
            return 0
        w = max(px(sh.width) - 40, 60)
        total = 0
        for para in sh.text_frame.paragraphs:
            text = "".join(r.text for r in para.runs)
            pts = [r.font.size.pt for r in para.runs if r.font.size]
            fpx = (max(pts) if pts else 12) * CW / 960.0
            if not text:
                total += fpx * 1.35
                continue
            cpl = max(int(w / (fpx * 0.46)), 8)    # ~0.46 em average advance
            total += max(1, math.ceil(len(text) / float(cpl))) * fpx * 1.35
        return total

    TOL = 24            # px of slack on the wrapped-text estimate
    fails = []
    for i, slide in enumerate(prs.slides, start=1):
        boxes = []
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            if x < -1 or y < -1 or x + w > CW + 1 or y + h > CH + 1:
                fails.append("slide %d: %s leaves the canvas (x=%.0f y=%.0f w=%.0f h=%.0f)"
                             % (i, sh.shape_type, x, y, w, h))
            th = est_text_h(sh)
            if th and y + th > CH + TOL:
                fails.append("slide %d: text overruns the bottom edge (y=%.0f + est %.0f > %d px) — \"%s\""
                             % (i, y, th, CH, sh.text_frame.text[:40].replace("\n", " ")))
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        if r.font.size and r.font.size.pt < FLOOR_PT - 0.25:
                            fails.append("slide %d: %.1f pt is below the %.1f pt floor (= %d px) — \"%s\""
                                         % (i, r.font.size.pt, FLOOR_PT, cfg.font_floor_px, r.text[:32]))
            boxes.append((sh, x, y, w, h))   # declared height: the estimate is
                                             # deliberately pessimistic and would
                                             # invent collisions with captions
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                sa, xa, ya, wa, ha = boxes[a]
                sb, xb, yb, wb, hb = boxes[b]
                if sa.shape_type != PICTURE and sb.shape_type != PICTURE:
                    continue                       # labels beside labels are fine
                if i == 1:
                    continue                       # title slide: text over the hero image is intended
                ov = min(xa + wa, xb + wb) - max(xa, xb)
                oh = min(ya + ha, yb + hb) - max(ya, yb)
                if ov > 2 and oh > 2:
                    fails.append("slide %d: overlap — %s over %s" % (i, sa.shape_type, sb.shape_type))
    for f in fails:
        print(f)
    print("PPTX LAYOUT %s (%d slides, %d finding%s)"
          % ("FAIL" if fails else "PASS", len(prs.slides._sldIdLst),
             len(fails), "" if len(fails) == 1 else "s"))
    return 1 if fails else 0


if __name__ == "__main__":
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    sys.exit(main())
